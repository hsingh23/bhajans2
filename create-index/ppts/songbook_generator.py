#!/usr/bin/env python3
"""
songbook_generator.py

Complete pipeline to convert PPTX song files to XML, PDF, and TOC.
Supports multiple year folders with robust error handling and detailed reporting.

Usage with uv:
    uv run songbook_generator.py /path/to/base/folder -o output
"""

import os
import sys
import re
import argparse
import unicodedata
import xml.etree.ElementTree as ET
from xml.dom import minidom
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional
from datetime import datetime
from enum import Enum, auto
import json
import traceback

from pptx import Presentation
from reportlab.lib.units import inch
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak,
    Flowable, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


# =============================================================================
# CONFIGURATION CONSTANTS
# =============================================================================

# ------------------------------------------------------------------------------
# YEAR FOLDERS
# ------------------------------------------------------------------------------
YEAR_FOLDERS = ["2021", "2022", "2023", "2024", "2025"]
CURRENT_YEAR = "2025"
DEFAULT_LANGUAGE = None  # Set to e.g., "Unknown" to have a fallback

# ------------------------------------------------------------------------------
# FONT CONFIGURATION
# ------------------------------------------------------------------------------
# For web PDFs with Unicode, we must embed fonts. 
# Uses font subsetting to minimize file size.
FONT_REGULAR_PATH = None
FONT_BOLD_PATH = None
FONT_ITALIC_PATH = None
FONT_BOLD_ITALIC_PATH = None
FONT_FAMILY = "SongFont"

# ------------------------------------------------------------------------------
# PDF PAGE SETTINGS
# ------------------------------------------------------------------------------
PAGE_WIDTH = 5.5 * inch
PAGE_HEIGHT = 8.5 * inch
MARGIN_LEFT = 0.5 * inch
MARGIN_RIGHT = 0.5 * inch
MARGIN_TOP = 0.7 * inch
MARGIN_BOTTOM = 0.6 * inch

# ------------------------------------------------------------------------------
# PDF HEADER SETTINGS
# ------------------------------------------------------------------------------
HEADER_TEXT = "Devotional Songs of Sri Mata Amritanandamayi"
HEADER_FONT_SIZE = 9
HEADER_Y_POSITION = 0.4 * inch
HEADER_LINE_ENABLED = True
HEADER_LINE_Y_POSITION = 0.55 * inch
HEADER_LINE_THICKNESS = 0.5

# ------------------------------------------------------------------------------
# PDF PAGE NUMBER
# ------------------------------------------------------------------------------
PAGE_NUMBER_FONT_SIZE = 10
PAGE_NUMBER_Y_POSITION = 0.35 * inch

# ------------------------------------------------------------------------------
# PDF TYPOGRAPHY - TITLE
# ------------------------------------------------------------------------------
TITLE_FONT_SIZE = 22
TITLE_LINE_HEIGHT = 26
TITLE_SPACE_BEFORE = 0
TITLE_SPACE_AFTER = 2
TITLE_UNDERLINE_ENABLED = True
TITLE_UNDERLINE_THICKNESS = 0.5
TITLE_UNDERLINE_COLOR = colors.black

# ------------------------------------------------------------------------------
# PDF TYPOGRAPHY - SEARCHABLE TITLE
# ------------------------------------------------------------------------------
# Invisible ASCII version of title for Ctrl+F search functionality
# This allows users to search "prema-sagarame" and find "prēma-sāgaramē"
SEARCHABLE_TITLE_ENABLED = True
SEARCHABLE_TITLE_FONT_SIZE = 1        # 1pt - effectively invisible
SEARCHABLE_TITLE_LINE_HEIGHT = 1
SEARCHABLE_TITLE_SPACE_AFTER = 2
# Transparent color (won't show but is selectable/searchable)
SEARCHABLE_TITLE_COLOR = colors.Color(1, 1, 1, alpha=0)
# Alternative: white (blends with white background)
# SEARCHABLE_TITLE_COLOR = colors.white
# Alternative: very light gray (barely visible)
# SEARCHABLE_TITLE_COLOR = colors.Color(0.95, 0.95, 0.95)

# ------------------------------------------------------------------------------
# PDF TYPOGRAPHY - LYRICS
# ------------------------------------------------------------------------------
LYRICS_FONT_SIZE = 13
LYRICS_LINE_HEIGHT = 18
LYRICS_SPACE_BEFORE = 16
LYRICS_SPACE_AFTER = 6

# ------------------------------------------------------------------------------
# PDF TYPOGRAPHY - MEANING
# ------------------------------------------------------------------------------
MEANING_FONT_SIZE = 11
MEANING_LINE_HEIGHT = 14
MEANING_SPACE_BEFORE = 4
MEANING_SPACE_AFTER = 20

# ------------------------------------------------------------------------------
# STANZA SPACING
# ------------------------------------------------------------------------------
STANZA_SEPARATOR_HEIGHT = 12

# ------------------------------------------------------------------------------
# TOC SETTINGS
# ------------------------------------------------------------------------------
TOC_FORMAT = "{title} ({language}) ## {year}-{page}"
TOC_FORMAT_NO_LANGUAGE = "{title} ## {year}-{page}"
TOC_SORT_ALPHABETICALLY = True

# ------------------------------------------------------------------------------
# REPORT SETTINGS
# ------------------------------------------------------------------------------
REPORT_EXTENSION = "_report.txt"
REPORT_JSON_EXTENSION = "_report.json"

# ------------------------------------------------------------------------------
# XML SETTINGS
# ------------------------------------------------------------------------------
XML_INDENT = "  "
XML_INCLUDE_SOURCE = True

# ------------------------------------------------------------------------------
# FILE NAMING
# ------------------------------------------------------------------------------
DEFAULT_OUTPUT_NAME = "songbook"
XML_EXTENSION = ".xml"
PDF_EXTENSION = ".pdf"
TOC_EXTENSION = "_toc.txt"

# ------------------------------------------------------------------------------
# KNOWN LANGUAGES AND DETECTION
# ------------------------------------------------------------------------------
# Full language names and their abbreviations/variations
# Format: canonical_name -> [list of variations/abbreviations]
LANGUAGE_MAP = {
    # Indian Languages
    "Malayalam": [
        "malayalam", "mal", "mlm", "malyalam", "malaylam", "malay",
        "ml", "മലയാളം"
    ],
    "Tamil": [
        "tamil", "tam", "tml", "tamizh", "thamil", "ta", "தமிழ்"
    ],
    "Telugu": [
        "telugu", "tel", "tlg", "telegu", "te", "తెలుగు"
    ],
    "Kannada": [
        "kannada", "kan", "knd", "kannad", "karnataka", "kn", "ಕನ್ನಡ"
    ],
    "Hindi": [
        "hindi", "hin", "hnd", "hdi", "hi", "हिंदी", "हिन्दी"
    ],
    "Sanskrit": [
        "sanskrit", "san", "skt", "samskrit", "samskritam", "sa", "संस्कृत"
    ],
    "Bengali": [
        "bengali", "ben", "bng", "bangla", "bangali", "bn", "বাংলা"
    ],
    "Marathi": [
        "marathi", "mar", "mrt", "marati", "mr", "मराठी"
    ],
    "Gujarati": [
        "gujarati", "guj", "gujrati", "gujarathi", "gu", "ગુજરાતી"
    ],
    "Punjabi": [
        "punjabi", "pun", "pnj", "panjabi", "pa", "ਪੰਜਾਬੀ", "پنجابی"
    ],
    "Odia": [
        "odia", "odi", "oriya", "orya", "odiya", "odhiya", "or", "ଓଡ଼ିଆ"
    ],
    "Braj Bhasha": [
        "braj bhasha", "braj", "braja", "brajabhasha", "brj"
    ],
    "Marwari": [
        "marwari", "marwadi", "mwr", "marvadi"
    ],
    "Assamese": [
        "assamese", "asm", "assam", "asamiya", "as", "অসমীয়া"
    ],
    "Konkani": [
        "konkani", "kon", "knk", "kk", "कोंকণী"
    ],
    "Kashmiri": [
        "kashmiri", "kas", "kash", "ks", "कॉशुर"
    ],
    "Nepali": [
        "nepali", "nep", "npl", "ne", "नेपाली"
    ],
    "Sindhi": [
        "sindhi", "snd", "sind", "sd", "سنڌي"
    ],
    "Maithili": [
        "maithili", "mai", "mth", "मैथिली"
    ],
    "Dogri": [
        "dogri", "doi", "dgr", "डोगरी"
    ],
    "Manipuri": [
        "manipuri", "mni", "mnp", "meitei", "মৈতৈলোন্"
    ],
    "Bodo": [
        "bodo", "bod", "brx", "बड़ो"
    ],
    "Santali": [
        "santali", "sat", "snt", "ᱥᱟᱱᱛᱟᱲᱤ"
    ],
    "Urdu": [
        "urdu", "urd", "ur", "اردو"
    ],
    "Tulu": [
        "tulu", "tcy", "ತುಳು"
    ],
    "Bhojpuri": [
        "bhojpuri", "bho", "bhoj", "भोजपुरी"
    ],
    "Rajasthani": [
        "rajasthani", "raj", "rjs", "राजस्थानी"
    ],
    "Chhattisgarhi": [
        "chhattisgarhi", "chg", "छत्तीसगढ़ी"
    ],
    "Haryanvi": [
        "haryanvi", "har", "hry", "हरियाणवी"
    ],
    "Magahi": [
        "magahi", "mag", "मगही"
    ],
    "Awadhi": [
        "awadhi", "awa", "अवधी"
    ],
    
    # Other Languages
    "English": [
        "english", "eng", "en"
    ],
    "Spanish": [
        "spanish", "spa", "español", "espanol", "es"
    ],
    "French": [
        "french", "fra", "fre", "français", "francais", "fr"
    ],
    "German": [
        "german", "ger", "deu", "deutsch", "de"
    ],
    "Portuguese": [
        "portuguese", "por", "português", "pt"
    ],
    "Italian": [
        "italian", "ita", "italiano", "it"
    ],
    "Russian": [
        "russian", "rus", "ru", "русский"
    ],
    "Japanese": [
        "japanese", "jpn", "ja", "日本語"
    ],
    "Chinese": [
        "chinese", "chi", "zho", "zh", "中文", "mandarin"
    ],
    "Korean": [
        "korean", "kor", "ko", "한국어"
    ],
    "Arabic": [
        "arabic", "ara", "ar", "العربية"
    ],
    "Persian": [
        "persian", "fas", "farsi", "fa", "فارسی"
    ],
    "Turkish": [
        "turkish", "tur", "tr", "türkçe"
    ],
    "Thai": [
        "thai", "tha", "th", "ไทย"
    ],
    "Vietnamese": [
        "vietnamese", "vie", "vi", "tiếng việt"
    ],
    "Indonesian": [
        "indonesian", "ind", "id", "bahasa"
    ],
    "Malay": [
        "malay", "msa", "ms", "melayu"
    ],
    "Swahili": [
        "swahili", "swa", "sw", "kiswahili"
    ],
    "Hebrew": [
        "hebrew", "heb", "he", "עברית"
    ],
}

# Build reverse lookup: variation -> canonical name
LANGUAGE_LOOKUP: dict[str, str] = {}
for canonical, variations in LANGUAGE_MAP.items():
    LANGUAGE_LOOKUP[canonical.lower()] = canonical
    for var in variations:
        LANGUAGE_LOOKUP[var.lower()] = canonical

# Common patterns in filenames that indicate language
# e.g., "Song Name 2025 Mal in Eng.pptx" or "Song_Tamil_2024.pptx"
FILENAME_LANGUAGE_PATTERNS = [
    # "Mal in Eng" pattern - source language
    r'\b([A-Za-z]{2,})\s+in\s+[Ee]ng(?:lish)?\b',
    # "in Malayalam" pattern
    r'\bin\s+([A-Za-z]{2,})\b',
    # Language as a standalone word
    r'\b([A-Za-z]{3,})\b',
    # Language code in parentheses: (Mal) or (Tamil)
    r'\(([A-Za-z]{2,})\)',
    # Language after year: "2025 Malayalam" or "2025_Tamil"
    r'20\d{2}[\s_\-]+([A-Za-z]{2,})',
    # Language before year: "Malayalam 2025"
    r'([A-Za-z]{3,})[\s_\-]+20\d{2}',
    # Underscored: "Song_Malayalam_2025"
    r'_([A-Za-z]{3,})_',
    # Hyphenated: "Song-Tamil-2025"
    r'-([A-Za-z]{3,})-',
]


# =============================================================================
# DATA CLASSES
# =============================================================================

class IssueType(Enum):
    """Types of issues during processing."""
    MISSING_TITLE = auto()
    MISSING_YEAR = auto()
    MISSING_LANGUAGE = auto()
    MISSING_MEANING = auto()
    EMPTY_STANZA = auto()
    NO_STANZAS = auto()
    PARSE_ERROR = auto()
    FILE_ERROR = auto()
    SLIDE_ERROR = auto()


@dataclass
class Issue:
    """Detailed issue record."""
    issue_type: IssueType
    source_file: str
    full_path: str = ""
    song_title: Optional[str] = None
    slide_number: Optional[int] = None
    stanza_number: Optional[int] = None
    details: str = ""
    raw_text: Optional[str] = None  # For debugging
    
    def to_dict(self) -> dict:
        return {
            'type': self.issue_type.name,
            'source_file': self.source_file,
            'full_path': self.full_path,
            'song_title': self.song_title,
            'slide_number': self.slide_number,
            'stanza_number': self.stanza_number,
            'details': self.details,
            'raw_text': self.raw_text,
        }


@dataclass
class Stanza:
    """A stanza with lyrics and optional meaning."""
    lyrics: str
    meaning: Optional[str] = None
    slide_number: int = 0
    has_meaning: bool = True


@dataclass 
class Song:
    """A complete song."""
    title: str
    year: Optional[str] = None
    language: Optional[str] = None
    source_file: Optional[str] = None
    full_path: Optional[str] = None
    stanzas: list[Stanza] = field(default_factory=list)
    page_number: int = 1
    
    # Source tracking
    year_source: str = ""      # "slide", "folder", "filename", "default"
    language_source: str = ""  # "slide_parens", "slide_text", "filename", ""
    has_year: bool = True
    has_language: bool = True
    missing_meaning_slides: list[int] = field(default_factory=list)


@dataclass
class ProcessingReport:
    """Comprehensive processing report."""
    timestamp: str = ""
    input_path: str = ""
    output_path: str = ""
    
    # Counts
    total_files_found: int = 0
    total_files_processed: int = 0
    total_songs_extracted: int = 0
    total_stanzas: int = 0
    total_slides_processed: int = 0
    
    # By year
    songs_by_year: dict = field(default_factory=dict)
    files_by_year: dict = field(default_factory=dict)
    
    # Issues
    issues: list[Issue] = field(default_factory=list)
    
    # Language detection (new)
    language_sources: dict[str, int] = field(default_factory=lambda: {
        "slide_parens": 0, "slide_text": 0, "filename": 0, "not_found": 0
    })
    
    # Aggregated counts
    songs_missing_year: list[str] = field(default_factory=list)
    songs_missing_language: list[str] = field(default_factory=list)
    songs_with_missing_meanings: dict = field(default_factory=dict)  # song -> list of slides
    files_with_errors: list[str] = field(default_factory=list)
    skipped_files: list[str] = field(default_factory=list)
    
    def add_issue(self, issue: Issue):
        self.issues.append(issue)
        
        key = f"{issue.source_file}"
        if issue.song_title:
            key = f"{issue.song_title} ({issue.source_file})"
        
        if issue.issue_type == IssueType.MISSING_YEAR:
            if key not in self.songs_missing_year:
                self.songs_missing_year.append(key)
        elif issue.issue_type == IssueType.MISSING_LANGUAGE:
            if key not in self.songs_missing_language:
                self.songs_missing_language.append(key)
        elif issue.issue_type == IssueType.MISSING_MEANING:
            if key not in self.songs_with_missing_meanings:
                self.songs_with_missing_meanings[key] = []
            if issue.slide_number:
                self.songs_with_missing_meanings[key].append(issue.slide_number)
        elif issue.issue_type in (IssueType.PARSE_ERROR, IssueType.FILE_ERROR, IssueType.SLIDE_ERROR):
            if issue.full_path not in self.files_with_errors:
                self.files_with_errors.append(issue.full_path or issue.source_file)


# Global report
_report: Optional[ProcessingReport] = None


def get_report() -> ProcessingReport:
    global _report
    if _report is None:
        _report = ProcessingReport()
    return _report


def reset_report():
    global _report
    _report = ProcessingReport()
    _report.timestamp = datetime.now().isoformat()


# =============================================================================
# FONT REGISTRATION
# =============================================================================

def find_system_fonts() -> dict[str, Optional[str]]:
    """Find Unicode fonts on the system."""
    font_paths = {'regular': None, 'bold': None, 'italic': None, 'bold_italic': None}
    
    # Check custom paths first
    if FONT_REGULAR_PATH and os.path.exists(FONT_REGULAR_PATH):
        font_paths['regular'] = FONT_REGULAR_PATH
        font_paths['bold'] = FONT_BOLD_PATH if FONT_BOLD_PATH and os.path.exists(FONT_BOLD_PATH) else FONT_REGULAR_PATH
        font_paths['italic'] = FONT_ITALIC_PATH if FONT_ITALIC_PATH and os.path.exists(FONT_ITALIC_PATH) else FONT_REGULAR_PATH
        font_paths['bold_italic'] = FONT_BOLD_ITALIC_PATH if FONT_BOLD_ITALIC_PATH and os.path.exists(FONT_BOLD_ITALIC_PATH) else font_paths['bold']
        return font_paths
    
    # Font families to search
    search_patterns = [
        # Noto Serif
        ("NotoSerif", [
            "fonts/NotoSerif-{}.ttf",
            "/usr/share/fonts/truetype/noto/NotoSerif-{}.ttf",
            "/usr/share/fonts/noto/NotoSerif-{}.ttf",
            "C:/Windows/Fonts/NotoSerif-{}.ttf",
            "/Library/Fonts/NotoSerif-{}.ttf",
            os.path.expanduser("~/.local/share/fonts/NotoSerif-{}.ttf"),
        ], ["Regular", "Bold", "Italic", "BoldItalic"]),
        
        # DejaVu Serif  
        ("DejaVuSerif", [
            "fonts/DejaVuSerif{}.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSerif{}.ttf",
            "/usr/share/fonts/dejavu/DejaVuSerif{}.ttf",
        ], ["", "-Bold", "-Italic", "-BoldItalic"]),
        
        # Liberation Serif
        ("LiberationSerif", [
            "fonts/LiberationSerif-{}.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSerif-{}.ttf",
        ], ["Regular", "Bold", "Italic", "BoldItalic"]),
        
        # FreeSerif
        ("FreeSerif", [
            "fonts/FreeSerif{}.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSerif{}.ttf",
        ], ["", "Bold", "Italic", "BoldItalic"]),
    ]
    
    for family_name, path_patterns, suffixes in search_patterns:
        for pattern in path_patterns:
            regular_path = pattern.format(suffixes[0])
            if os.path.exists(regular_path):
                font_paths['regular'] = regular_path
                
                bold_path = pattern.format(suffixes[1])
                font_paths['bold'] = bold_path if os.path.exists(bold_path) else regular_path
                
                italic_path = pattern.format(suffixes[2])
                font_paths['italic'] = italic_path if os.path.exists(italic_path) else regular_path
                
                bold_italic_path = pattern.format(suffixes[3])
                font_paths['bold_italic'] = bold_italic_path if os.path.exists(bold_italic_path) else font_paths['bold']
                
                print(f"  Found font: {family_name}")
                return font_paths
    
    return font_paths


def register_fonts() -> tuple[str, str, str, str]:
    """Register fonts with ReportLab."""
    font_paths = find_system_fonts()
    
    if not font_paths['regular']:
        print("\n" + "!" * 60)
        print("WARNING: No Unicode fonts found!")
        print("!" * 60)
        print("Run: uv run download_fonts.py")
        print("!" * 60 + "\n")
        return ("Times-Roman", "Times-Bold", "Times-Italic", "Times-BoldItalic")
    
    try:
        font_regular = f"{FONT_FAMILY}-Regular"
        font_bold = f"{FONT_FAMILY}-Bold"
        font_italic = f"{FONT_FAMILY}-Italic"
        font_bold_italic = f"{FONT_FAMILY}-BoldItalic"
        
        pdfmetrics.registerFont(TTFont(font_regular, font_paths['regular']))
        pdfmetrics.registerFont(TTFont(font_bold, font_paths['bold']))
        pdfmetrics.registerFont(TTFont(font_italic, font_paths['italic']))
        pdfmetrics.registerFont(TTFont(font_bold_italic, font_paths['bold_italic']))
        
        from reportlab.pdfbase.pdfmetrics import registerFontFamily
        registerFontFamily(FONT_FAMILY, normal=font_regular, bold=font_bold, 
                          italic=font_italic, boldItalic=font_bold_italic)
        
        return (font_regular, font_bold, font_italic, font_bold_italic)
    except Exception as e:
        print(f"  Font registration error: {e}")
        return ("Times-Roman", "Times-Bold", "Times-Italic", "Times-BoldItalic")


# =============================================================================
# TEXT EXTRACTION
# =============================================================================

def extract_text_from_shape(shape) -> str:
    """
    Extract text from PowerPoint shape, preserving ALL line breaks.
    
    PowerPoint stores text in paragraphs. Each paragraph is a separate line.
    Within a paragraph, there might also be soft line breaks (Shift+Enter).
    """
    if not shape.has_text_frame:
        return ""
    
    paragraphs = []
    
    for paragraph in shape.text_frame.paragraphs:
        # Collect text from all runs in this paragraph
        para_text_parts = []
        
        for run in paragraph.runs:
            run_text = run.text
            # Check for soft line breaks (vertical tab character in PPTX)
            # These are represented as '\v' or '\x0b' in Python
            run_text = run_text.replace('\v', '\n').replace('\x0b', '\n')
            para_text_parts.append(run_text)
        
        para_text = ''.join(para_text_parts)
        paragraphs.append(para_text)
    
    # Join paragraphs with newlines
    result = '\n'.join(paragraphs)
    
    # Clean up: remove leading/trailing blank lines but preserve internal structure
    lines = result.split('\n')
    
    # Strip leading empty lines
    while lines and not lines[0].strip():
        lines.pop(0)
    
    # Strip trailing empty lines  
    while lines and not lines[-1].strip():
        lines.pop()
    
    return '\n'.join(lines)


def get_shapes_by_position(slide) -> list[dict]:
    """Get text shapes sorted by vertical then horizontal position."""
    shapes = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text.strip():
                shapes.append({
                    'top': shape.top or 0,
                    'left': shape.left or 0,
                    'text': text,
                    'raw': text,  # Keep original for debugging
                })
    
    shapes.sort(key=lambda x: (x['top'], x['left']))
    return shapes


# =============================================================================
# LANGUAGE DETECTION
# =============================================================================

def normalize_language(text: str) -> Optional[str]:
    """
    Normalize a language string to its canonical form.
    
    Examples:
        "mal" -> "Malayalam"
        "Tam" -> "Tamil"
        "KANNADA" -> "Kannada"
        "sanskrit" -> "Sanskrit"
    """
    if not text:
        return None
    
    clean = text.strip().lower()
    
    # Direct lookup
    if clean in LANGUAGE_LOOKUP:
        return LANGUAGE_LOOKUP[clean]
    
    # Try without common suffixes/prefixes
    for suffix in ['i', 'am', 'u']:
        if clean.endswith(suffix) and clean[:-len(suffix)] in LANGUAGE_LOOKUP:
            return LANGUAGE_LOOKUP[clean[:-len(suffix)]]
    
    return None


def detect_language_from_text(text: str) -> Optional[str]:
    """
    Detect language from arbitrary text (slide content, metadata, etc.)
    
    Handles formats like:
        - "Malayalam"
        - "(Malayalam)"
        - "2025 (Malayalam)"
        - "2025 Malayalam"
        - "(mal)"
        - "Mal"
    """
    if not text:
        return None
    
    text = text.strip()
    
    # Check for language in parentheses first (highest priority)
    paren_matches = re.findall(r'\(([^)]+)\)', text)
    for match in paren_matches:
        lang = normalize_language(match.strip())
        if lang:
            return lang
    
    # Check for "in <language>" pattern
    in_match = re.search(r'\bin\s+([A-Za-z]+)', text, re.IGNORECASE)
    if in_match:
        lang = normalize_language(in_match.group(1))
        if lang and lang != "English":  # Skip "in English" as that's the translation
            return lang
    
    # Check for language after year
    after_year = re.search(r'20\d{2}\s+([A-Za-z]+)', text)
    if after_year:
        lang = normalize_language(after_year.group(1))
        if lang:
            return lang
    
    # Check if entire text (cleaned) is a language
    clean_text = re.sub(r'[^A-Za-z]', '', text)
    if clean_text:
        lang = normalize_language(clean_text)
        if lang:
            return lang
    
    # Check each word
    words = re.findall(r'[A-Za-z]+', text)
    for word in words:
        lang = normalize_language(word)
        if lang:
            return lang
    
    return None


def detect_language_from_filename(filename: str) -> Optional[str]:
    """
    Extract language from filename.
    
    Handles patterns like:
        - "Prema Sagarame 2025 Mal in Eng.pptx" -> Malayalam
        - "Bhajan_Tamil_2024.pptx" -> Tamil
        - "Song (Malayalam) 2025.pptx" -> Malayalam
        - "2025 Kannada - Song Name.pptx" -> Kannada
        - "Song Name Telugu.pptx" -> Telugu
    """
    if not filename:
        return None
    
    # Remove extension
    name = os.path.splitext(filename)[0]
    
    # Try each pattern
    for pattern in FILENAME_LANGUAGE_PATTERNS:
        matches = re.findall(pattern, name, re.IGNORECASE)
        for match in matches:
            lang = normalize_language(match)
            if lang:
                # Skip "Eng" or "English" as that usually indicates translation target
                if lang != "English":
                    return lang
    
    # Fallback: check all words in filename
    words = re.findall(r'[A-Za-z]+', name)
    for word in words:
        # Skip common non-language words
        skip_words = {
            'in', 'the', 'and', 'of', 'for', 'with', 'eng', 'english',
            'song', 'bhajan', 'kirtan', 'stotram', 'mantra', 'prayer',
            'devotional', 'spiritual', 'vol', 'volume', 'part', 'new',
            'pptx', 'ppt', 'pdf', 'mp3', 'mp4', 'lyrics', 'meaning',
            'translation', 'trans', 'final', 'draft', 'copy', 'old',
        }
        if word.lower() in skip_words:
            continue
        
        lang = normalize_language(word)
        if lang and lang != "English":
            return lang
    
    return None


def detect_language(
    slide_text: Optional[str] = None,
    filename: Optional[str] = None,
    all_slide_texts: Optional[list[str]] = None
) -> tuple[Optional[str], str]:
    """
    Detect language using multiple sources.
    
    Priority:
        1. Explicit language on title slide (in parentheses)
        2. Language in filename
        3. Language mentioned anywhere on title slide
        4. Language from other slides (less reliable)
    
    Args:
        slide_text: Text from title slide
        filename: The PPTX filename
        all_slide_texts: Optional list of text from all shapes on title slide
    
    Returns:
        (language, source) where source is one of:
        "slide_parens", "slide_text", "filename", or ""
    """
    # 1. Check for language in parentheses on slide (most explicit)
    if slide_text:
        paren_matches = re.findall(r'\(([^)]+)\)', slide_text)
        for match in paren_matches:
            lang = normalize_language(match.strip())
            if lang:
                return lang, "slide_parens"
    
    # 2. Check slide text more broadly
    if slide_text:
        lang = detect_language_from_text(slide_text)
        if lang:
            return lang, "slide_text"
    
    # 3. Check all slide texts
    if all_slide_texts:
        for text in all_slide_texts:
            # First check parentheses
            paren_matches = re.findall(r'\(([^)]+)\)', text)
            for match in paren_matches:
                lang = normalize_language(match.strip())
                if lang:
                    return lang, "slide_parens"
        
        # Then check text content
        for text in all_slide_texts:
            lang = detect_language_from_text(text)
            if lang:
                return lang, "slide_text"
    
    # 4. Check filename
    if filename:
        lang = detect_language_from_filename(filename)
        if lang:
            return lang, "filename"
    
    return None, ""


def detect_year(text: str) -> Optional[str]:
    """Extract year (4 digits between 2000-2099) from text."""
    match = re.search(r'(20[0-9]{2})', text)
    return match.group(1) if match else None


# =============================================================================
# PPTX PARSING
# =============================================================================

def parse_title_slide(
    slide, 
    source_file: str, 
    full_path: str,
    folder_year: Optional[str] = None
) -> tuple[Optional[str], Optional[str], Optional[str], str, str]:
    """
    Parse title slide for title, year, and language.
    
    Returns:
        (title, year, language, year_source, language_source)
        year_source: "slide", "folder", "filename", "default"
        language_source: "slide_parens", "slide_text", "filename", ""
    """
    report = get_report()
    shapes = get_shapes_by_position(slide)
    
    if not shapes:
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_TITLE,
            source_file=source_file,
            full_path=full_path,
            slide_number=1,
            details="No text found on title slide (slide 1)",
        ))
        return None, None, None, "", ""
    
    # Collect all text for analysis
    all_texts = [s['text'] for s in shapes]
    all_text_combined = "\n---\n".join(all_texts)
    
    # First shape is typically the title
    title_text = shapes[0]['text'].strip()
    
    # Handle multi-line titles
    if '\n' in title_text:
        lines = [l.strip() for l in title_text.split('\n') if l.strip()]
        title = lines[0] if lines else ""
    else:
        title = title_text
    
    if not title:
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_TITLE,
            source_file=source_file,
            full_path=full_path,
            slide_number=1,
            details="Title text box is empty",
            raw_text=all_text_combined[:500],
        ))
        return None, None, None, "", ""
    
    # --- YEAR DETECTION ---
    year = None
    year_source = ""
    
    # Check all slide text for year
    for text in all_texts:
        found_year = detect_year(text)
        if found_year:
            year = found_year
            year_source = "slide"
            break
    
    # Try filename for year
    if not year:
        found_year = detect_year(source_file)
        if found_year:
            year = found_year
            year_source = "filename"
    
    # Fallback to folder year
    if not year and folder_year:
        year = folder_year
        year_source = "folder"
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_YEAR,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            slide_number=1,
            details=f"Year not found on slide or filename, using folder: {folder_year}",
            raw_text=all_text_combined[:300],
        ))
    
    # Final fallback to default
    if not year:
        year = CURRENT_YEAR
        year_source = "default"
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_YEAR,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            slide_number=1,
            details=f"Year not found anywhere, using default: {CURRENT_YEAR}",
            raw_text=all_text_combined[:300],
        ))
    
    # --- LANGUAGE DETECTION ---
    language, language_source = detect_language(
        slide_text=all_text_combined,
        filename=source_file,
        all_slide_texts=all_texts
    )
    
    if not language:
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_LANGUAGE,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            slide_number=1,
            details=(
                "Language not found. Checked:\n"
                f"  - Title slide text\n"
                f"  - Filename: {source_file}\n"
                "Expected formats: '(Malayalam)', '2025 Tamil', 'Song_Mal_2025.pptx'"
            ),
            raw_text=all_text_combined[:400],
        ))
    
    return title, year, language, year_source, language_source


def parse_lyric_slide(
    slide,
    slide_number: int,
    source_file: str,
    full_path: str,
    song_title: str,
    language: Optional[str] = None
) -> tuple[Optional[str], Optional[str], bool]:
    """
    Parse lyric slide.
    
    Returns:
        (lyrics, meaning, had_error)
    """
    report = get_report()
    
    try:
        shapes = get_shapes_by_position(slide)
    except Exception as e:
        report.add_issue(Issue(
            issue_type=IssueType.SLIDE_ERROR,
            source_file=source_file,
            full_path=full_path,
            song_title=song_title,
            slide_number=slide_number,
            details=f"Error reading slide: {str(e)}",
        ))
        return None, None, True
    
    if not shapes:
        report.add_issue(Issue(
            issue_type=IssueType.EMPTY_STANZA,
            source_file=source_file,
            full_path=full_path,
            song_title=song_title,
            slide_number=slide_number,
            details=f"Slide {slide_number} has no text content",
        ))
        return None, None, False
    
    if len(shapes) == 1:
        # Only lyrics, no meaning
        lyrics = shapes[0]['text']
        
        # English songs don't need translation/meaning
        if language != "English":
            report.add_issue(Issue(
                issue_type=IssueType.MISSING_MEANING,
                source_file=source_file,
                full_path=full_path,
                song_title=song_title,
                slide_number=slide_number,
                details=f"Slide {slide_number} has lyrics but no English meaning (only 1 text box found)",
                raw_text=lyrics[:200] if lyrics else None,
            ))
        return lyrics, None, False
    
    # First shape = lyrics, last shape = meaning
    lyrics = shapes[0]['text']
    meaning = shapes[-1]['text']
    
    # Sanity check: if meaning looks like more lyrics (no English), flag it
    # This is a heuristic - you can adjust or remove
    if language != "English" and meaning and not any(c.isascii() and c.isalpha() for c in meaning[:50]):
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_MEANING,
            source_file=source_file,
            full_path=full_path,
            song_title=song_title,
            slide_number=slide_number,
            details=f"Slide {slide_number}: Bottom text doesn't appear to be English translation",
            raw_text=f"Top: {lyrics[:100]}...\nBottom: {meaning[:100]}...",
        ))
        # Still return it - let user decide
    
    return lyrics, meaning, False


def process_pptx_file(
    filepath: str, 
    folder_year: Optional[str] = None
) -> Optional[Song]:
    """
    Process a single PPTX file.
    """
    report = get_report()
    source_file = os.path.basename(filepath)
    full_path = os.path.abspath(filepath)
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        report.add_issue(Issue(
            issue_type=IssueType.FILE_ERROR,
            source_file=source_file,
            full_path=full_path,
            details=f"Cannot open file: {str(e)}\n{traceback.format_exc()}",
        ))
        return None
    
    slides = list(prs.slides)
    if not slides:
        report.add_issue(Issue(
            issue_type=IssueType.PARSE_ERROR,
            source_file=source_file,
            full_path=full_path,
            details="File contains no slides",
        ))
        return None
    
    report.total_slides_processed += len(slides)
    
    # Parse title slide
    title, year, language, year_source, language_source = parse_title_slide(
        slides[0], source_file, full_path, folder_year
    )
    
    if not title:
        return None
    
    # Parse lyric slides
    stanzas = []
    missing_meaning_slides = []
    
    for idx, slide in enumerate(slides[1:], start=2):
        lyrics, meaning, had_error = parse_lyric_slide(
            slide, idx, source_file, full_path, title, language
        )
        
        if had_error:
            continue
        
        if not lyrics:
            continue
        
        has_meaning = meaning is not None and meaning.strip() != ""
        
        if not has_meaning and language != "English":
            missing_meaning_slides.append(idx)
        
        stanzas.append(Stanza(
            lyrics=lyrics,
            meaning=meaning if has_meaning else None,
            slide_number=idx,
            has_meaning=has_meaning,
        ))
    
    if not stanzas:
        report.add_issue(Issue(
            issue_type=IssueType.NO_STANZAS,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            details=f"No valid lyric slides found. File has {len(slides)} slides total.",
        ))
    
    return Song(
        title=title,
        year=year,
        language=language,
        source_file=source_file,
        full_path=full_path,
        stanzas=stanzas,
        year_source=year_source,
        language_source=language_source,
        has_year=(year_source == "slide"),
        has_language=(language is not None),
        missing_meaning_slides=missing_meaning_slides,
    )


def process_year_folder(folder_path: str, year: str) -> list[Song]:
    """Process all PPTX in a year folder."""
    report = get_report()
    songs = []
    folder = Path(folder_path)
    
    if not folder.exists():
        print(f"  ⚠ Folder not found: {folder_path}")
        return songs
    
    pptx_files = [f for f in sorted(folder.glob("*.pptx")) if not f.name.startswith("~$")]
    
    if not pptx_files:
        print(f"  ⚠ No .pptx files in {folder_path}")
        return songs
    
    report.total_files_found += len(pptx_files)
    report.files_by_year[year] = len(pptx_files)
    
    print(f"  Found {len(pptx_files)} files")
    
    for filepath in pptx_files:
        song = process_pptx_file(str(filepath), folder_year=year)
        
        if song:
            songs.append(song)
            report.total_files_processed += 1
            report.total_songs_extracted += 1
            report.total_stanzas += len(song.stanzas)
            
            y = song.year or "Unknown"
            report.songs_by_year[y] = report.songs_by_year.get(y, 0) + 1
            
            # Update language source stats
            if not song.has_language:
                report.language_sources["not_found"] += 1
            else:
                src = song.language_source or "slide_text" # Fallback
                report.language_sources[src] = report.language_sources.get(src, 0) + 1
            
            # Detailed status
            issues = []
            if not song.has_language:
                issues.append("no language detected")
            elif song.language_source == "filename":
                issues.append(f"language from filename: {song.language}")
            
            if song.year_source == "folder":
                issues.append(f"year from folder: {song.year}")
            elif song.year_source == "filename":
                issues.append(f"year from filename: {song.year}")
            elif song.year_source == "default":
                issues.append(f"year defaulted: {song.year}")
            
            if song.missing_meaning_slides:
                issues.append(f"missing meaning: slides {song.missing_meaning_slides}")
            
            # Display
            lang_display = f" ({song.language})" if song.language else ""
            
            if issues:
                print(f"    ⚠ {song.title}{lang_display}")
                for issue in issues:
                    print(f"        └─ {issue}")
            else:
                print(f"    ✓ {song.title}{lang_display}")
        else:
            report.skipped_files.append(str(filepath))
            print(f"    ✗ {filepath.name}")
            print(f"        └─ skipped (see report for details)")
    
    return songs


def process_all_years(base_path: str, years: list[str] = None) -> list[Song]:
    """Process all year folders."""
    if years is None:
        years = YEAR_FOLDERS
    
    all_songs = []
    base = Path(base_path)
    
    for year in years:
        year_folder = base / year
        print(f"\n{'─' * 50}")
        print(f"Processing {year}/")
        print(f"{'─' * 50}")
        
        songs = process_year_folder(str(year_folder), year)
        all_songs.extend(songs)
        
        print(f"  → {len(songs)} songs from {year}")
    
    return all_songs


# =============================================================================
# XML GENERATION
# =============================================================================

def songs_to_xml(songs: list[Song], output_path: str) -> None:
    """Create XML from songs."""
    root = ET.Element('songbook')
    root.set('generated', datetime.now().isoformat())
    root.set('total_songs', str(len(songs)))
    
    for song in songs:
        song_elem = ET.SubElement(root, 'song')
        
        ET.SubElement(song_elem, 'title').text = song.title
        
        if song.year:
            year_elem = ET.SubElement(song_elem, 'year')
            year_elem.text = song.year
            year_elem.set('source', song.year_source)
        
        if song.language:
            ET.SubElement(song_elem, 'language').text = song.language
        
        if song.source_file:
            ET.SubElement(song_elem, 'source').text = song.source_file
        
        if song.full_path:
            ET.SubElement(song_elem, 'full_path').text = song.full_path
        
        # Metadata
        meta = ET.SubElement(song_elem, 'metadata')
        meta.set('has_year_on_slide', str(song.has_year).lower())
        meta.set('has_language', str(song.has_language).lower())
        meta.set('year_source', song.year_source)
        if song.missing_meaning_slides:
            meta.set('missing_meaning_slides', ','.join(map(str, song.missing_meaning_slides)))
        
        # Stanzas
        if song.stanzas:
            stanzas_elem = ET.SubElement(song_elem, 'stanzas')
            for stanza in song.stanzas:
                stanza_elem = ET.SubElement(stanzas_elem, 'stanza')
                stanza_elem.set('slide', str(stanza.slide_number))
                stanza_elem.set('has_meaning', str(stanza.has_meaning).lower())
                
                ET.SubElement(stanza_elem, 'lyrics').text = stanza.lyrics
                
                if stanza.meaning:
                    ET.SubElement(stanza_elem, 'meaning').text = stanza.meaning
    
    xml_str = minidom.parseString(ET.tostring(root, encoding='unicode')).toprettyxml(indent=XML_INDENT)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(xml_str)
    
    print(f"XML saved: {output_path}")


def xml_to_songs(xml_path: str) -> list[Song]:
    """Load songs from XML."""
    tree = ET.parse(xml_path)
    root = tree.getroot()
    songs = []
    
    for song_elem in root.findall('song'):
        stanzas = []
        stanzas_elem = song_elem.find('stanzas')
        
        if stanzas_elem is not None:
            for stanza_elem in stanzas_elem.findall('stanza'):
                stanzas.append(Stanza(
                    lyrics=stanza_elem.findtext('lyrics', ''),
                    meaning=stanza_elem.findtext('meaning'),
                    slide_number=int(stanza_elem.get('slide', 0)),
                    has_meaning=stanza_elem.get('has_meaning', 'true').lower() == 'true',
                ))
        
        meta = song_elem.find('metadata')
        year_elem = song_elem.find('year')
        
        missing_slides = []
        if meta is not None and meta.get('missing_meaning_slides'):
            missing_slides = [int(x) for x in meta.get('missing_meaning_slides').split(',')]
        
        songs.append(Song(
            title=song_elem.findtext('title', ''),
            year=song_elem.findtext('year'),
            language=song_elem.findtext('language'),
            source_file=song_elem.findtext('source'),
            full_path=song_elem.findtext('full_path'),
            stanzas=stanzas,
            year_source=year_elem.get('source', '') if year_elem is not None else '',
            has_year=meta.get('has_year_on_slide', 'true').lower() == 'true' if meta is not None else True,
            has_language=meta.get('has_language', 'true').lower() == 'true' if meta is not None else True,
            missing_meaning_slides=missing_slides,
        ))
    
    return songs


# =============================================================================
# PDF GENERATION
# =============================================================================

def normalize_for_search(text: str) -> str:
    """
    Normalize text to ASCII for searchability.
    Removes diacritics, keeps basic punctuation.
    
    Examples:
        "prēma-sāgaramē" -> "prema-sagarame"
        "Ādhyātmata" -> "Adhyatmata"
        "hṛdayēśvarī" -> "hrdayesvari"
    """
    if not text:
        return ""
    
    # NFD decomposition separates base characters from combining diacritics
    decomposed = unicodedata.normalize('NFD', text)
    
    # Remove combining diacritical marks (category 'Mn')
    ascii_chars = []
    for char in decomposed:
        if unicodedata.category(char) == 'Mn':
            continue  # Skip diacritics
        # Handle special characters
        if char == 'ś' or char == 'ṣ':
            ascii_chars.append('s')
        elif char == 'ṇ' or char == 'ñ':
            ascii_chars.append('n')
        elif char == 'ṁ' or char == 'ṃ':
            ascii_chars.append('m')
        elif char == 'ṛ' or char == 'ṝ':
            ascii_chars.append('r')
        elif char == 'ḷ' or char == 'ḹ':
            ascii_chars.append('l')
        elif char == 'ṭ' or char == 'ṭh':
            ascii_chars.append('t')
        elif char == 'ḍ' or char == 'ḍh':
            ascii_chars.append('d')
        else:
            ascii_chars.append(char)
    
    result = ''.join(ascii_chars)
    
    # Normalize whitespace
    result = re.sub(r'\s+', ' ', result).strip()
    
    return result


class PageTracker(Flowable):
    """Track page numbers for TOC."""
    _registry: dict[str, int] = {}
    
    def __init__(self, key: str):
        Flowable.__init__(self)
        self.key = key
        self.width = 0
        self.height = 0
    
    def draw(self):
        PageTracker._registry[self.key] = self.canv.getPageNumber()
    
    @classmethod
    def get_page(cls, key: str) -> Optional[int]:
        return cls._registry.get(key)
    
    @classmethod
    def reset(cls):
        cls._registry = {}


class HorizontalRule(Flowable):
    """Horizontal line."""
    def __init__(self, width: float):
        Flowable.__init__(self)
        self.line_width = width
        self.height = TITLE_UNDERLINE_THICKNESS + 4
    
    def draw(self):
        self.canv.setStrokeColor(TITLE_UNDERLINE_COLOR)
        self.canv.setLineWidth(TITLE_UNDERLINE_THICKNESS)
        self.canv.line(0, 2, self.line_width, 2)


def create_styles(fonts: tuple[str, str, str, str]) -> dict:
    """Create PDF styles including searchable title."""
    regular, bold, italic, bold_italic = fonts
    
    return {
        'title': ParagraphStyle(
            'Title',
            fontSize=TITLE_FONT_SIZE,
            fontName=bold,
            leading=TITLE_LINE_HEIGHT,
            spaceBefore=TITLE_SPACE_BEFORE,
            spaceAfter=0,  # We'll add our own spacing after searchable title
        ),
        'searchable_title': ParagraphStyle(
            'SearchableTitle',
            fontSize=SEARCHABLE_TITLE_FONT_SIZE,
            fontName=regular,
            leading=SEARCHABLE_TITLE_FONT_SIZE,
            spaceBefore=0,
            spaceAfter=SEARCHABLE_TITLE_SPACE_AFTER,
            textColor=SEARCHABLE_TITLE_COLOR,
        ),
        'lyrics': ParagraphStyle(
            'Lyrics',
            fontSize=LYRICS_FONT_SIZE,
            fontName=regular,
            leading=LYRICS_LINE_HEIGHT,
            spaceBefore=LYRICS_SPACE_BEFORE,
            spaceAfter=LYRICS_SPACE_AFTER,
        ),
        'meaning': ParagraphStyle(
            'Meaning',
            fontSize=MEANING_FONT_SIZE,
            fontName=bold_italic,
            leading=MEANING_LINE_HEIGHT,
            spaceBefore=MEANING_SPACE_BEFORE,
            spaceAfter=MEANING_SPACE_AFTER,
        ),
    }


def draw_page(canvas, doc, fonts):
    """Draw header and page number."""
    regular, bold, italic, bold_italic = fonts
    canvas.saveState()
    
    # Header
    canvas.setFont(italic, HEADER_FONT_SIZE)
    canvas.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT - HEADER_Y_POSITION, HEADER_TEXT)
    
    if HEADER_LINE_ENABLED:
        canvas.setStrokeColor(colors.black)
        canvas.setLineWidth(HEADER_LINE_THICKNESS)
        canvas.line(MARGIN_LEFT, PAGE_HEIGHT - HEADER_LINE_Y_POSITION,
                   PAGE_WIDTH - MARGIN_RIGHT, PAGE_HEIGHT - HEADER_LINE_Y_POSITION)
    
    # Page number
    canvas.setFont(regular, PAGE_NUMBER_FONT_SIZE)
    canvas.drawCentredString(PAGE_WIDTH / 2, PAGE_NUMBER_Y_POSITION, str(canvas.getPageNumber()))
    
    canvas.restoreState()


def escape_xml(text: str) -> str:
    """Escape XML special characters."""
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def text_to_html(text: str) -> str:
    """Convert text to HTML, preserving line breaks."""
    return escape_xml(text).replace('\n', '<br/>')


def generate_pdf(songs: list[Song], output_path: str) -> list[Song]:
    """
    Generate PDF with:
    - Preserved line breaks in lyrics
    - Searchable ASCII titles (invisible text for Ctrl+F)
    - Proper Unicode rendering
    """
    print("  Registering fonts...")
    fonts = register_fonts()
    
    PageTracker.reset()
    styles = create_styles(fonts)
    content_width = PAGE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    
    story = []
    
    for idx, song in enumerate(songs):
        # === PAGE TRACKING ===
        story.append(PageTracker(f"song_{idx}"))
        
        # === TITLE ===
        # Display title with Unicode characters
        title_display = escape_xml(song.title)
        if song.language:
            title_display += f" ({escape_xml(song.language)})"
        
        story.append(Paragraph(title_display, styles['title']))
        
        # === SEARCHABLE TITLE (invisible) ===
        # ASCII-only version for Ctrl+F searching
        # e.g., "prēma-sāgaramē" becomes "prema-sagarame"
        ascii_title = normalize_for_search(song.title)
        if song.language:
            ascii_title += f" {normalize_for_search(song.language)}"
        
        story.append(Paragraph(ascii_title, styles['searchable_title']))
        
        # === TITLE UNDERLINE ===
        if TITLE_UNDERLINE_ENABLED:
            story.append(HorizontalRule(content_width))
        
        story.append(Spacer(1, 8))
        
        # === STANZAS ===
        for stanza_idx, stanza in enumerate(song.stanzas):
            stanza_parts = []
            
            # --- LYRICS ---
            if stanza.lyrics:
                # Use <br/> for line breaks - this is the most reliable method
                lyrics_html = text_to_html(stanza.lyrics)
                stanza_parts.append(Paragraph(lyrics_html, styles['lyrics']))
            
            # --- MEANING ---
            if stanza.meaning:
                # For English songs, the "meaning" is often just a duplicate 
                # or redundant. If language is English, we only show it if 
                # it's clearly different from lyrics (rare in this dataset)
                # or if it exists and we want to preserve it.
                # USER said: English songs don't need translation.
                if song.language != "English":
                    meaning_html = text_to_html(stanza.meaning)
                    stanza_parts.append(Paragraph(meaning_html, styles['meaning']))
                else:
                    # Optional: only show meaning for English if it looks like a 
                    # specific variant/meaning rather than just being there.
                    # For now, following USER'S intent to not have translation.
                    pass
            
            # Keep lyrics + meaning together when possible
            if len(stanza_parts) > 1:
                story.append(KeepTogether(stanza_parts))
            elif stanza_parts:
                story.extend(stanza_parts)
            
            # Space between stanzas
            if stanza_idx < len(song.stanzas) - 1:
                story.append(Spacer(1, STANZA_SEPARATOR_HEIGHT))
        
        # === PAGE BREAK ===
        if idx < len(songs) - 1:
            story.append(PageBreak())
    
    # === BUILD DOCUMENT ===
    doc = SimpleDocTemplate(
        output_path,
        pagesize=(PAGE_WIDTH, PAGE_HEIGHT),
        leftMargin=MARGIN_LEFT,
        rightMargin=MARGIN_RIGHT,
        topMargin=MARGIN_TOP,
        bottomMargin=MARGIN_BOTTOM,
    )
    
    def page_handler(canvas, doc):
        draw_page(canvas, doc, fonts)
    
    doc.build(story, onFirstPage=page_handler, onLaterPages=page_handler)
    
    # Update page numbers for TOC
    for idx, song in enumerate(songs):
        song.page_number = PageTracker.get_page(f"song_{idx}") or (idx + 1)
    
    print(f"PDF saved: {output_path}")
    print(f"  - {len(songs)} songs")
    print(f"  - Searchable ASCII titles included")
    
    return songs


# =============================================================================
# TOC GENERATION
# =============================================================================

def normalize_toc(text: str) -> str:
    """Normalize to ASCII."""
    decomposed = unicodedata.normalize('NFD', text)
    ascii_chars = [c for c in decomposed if unicodedata.category(c) != 'Mn']
    result = re.sub(r'[^A-Za-z0-9 \-\(\)]', '', ''.join(ascii_chars))
    return re.sub(r'\s+', ' ', result).strip()


def generate_toc(songs: list[Song], output_path: str):
    """Generate TOC file."""
    entries = []
    
    for song in songs:
        title = normalize_toc(song.title)
        lang = (song.language or '').lower()
        year = song.year or CURRENT_YEAR
        
        if lang:
            entry = TOC_FORMAT.format(title=title, language=lang, year=year, page=song.page_number)
        else:
            entry = TOC_FORMAT_NO_LANGUAGE.format(title=title, year=year, page=song.page_number)
        entries.append(entry)
    
    if TOC_SORT_ALPHABETICALLY:
        entries.sort(key=str.lower)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(entries))
    
    print(f"TOC saved: {output_path}")


# =============================================================================
# DETAILED REPORT
# =============================================================================

def generate_report(output_path: str, json_path: str):
    """Generate detailed report."""
    report = get_report()
    
    lines = [
        "=" * 80,
        "SONGBOOK PROCESSING REPORT",
        "=" * 80,
        f"Generated: {report.timestamp}",
        f"Input:     {report.input_path}",
        f"Output:    {report.output_path}",
        "",
        "─" * 80,
        "SUMMARY",
        "─" * 80,
        f"  Files found:              {report.total_files_found}",
        f"  Files processed:          {report.total_files_processed}",
        f"  Files skipped/errored:    {len(report.skipped_files) + len(report.files_with_errors)}",
        f"  Songs extracted:          {report.total_songs_extracted}",
        f"  Total stanzas:            {report.total_stanzas}",
        f"  Total slides processed:   {report.total_slides_processed}",
        "",
    ]
    
    # Songs by year
    lines.extend([
        "─" * 80,
        "SONGS BY YEAR",
        "─" * 80,
    ])
    for year in sorted(report.songs_by_year.keys()):
        files_count = report.files_by_year.get(year, "?")
        lines.append(f"  {year}: {report.songs_by_year[year]} songs (from {files_count} files)")
    lines.append("")
    
    # Language detection summary
    lines.extend([
        "─" * 80,
        "LANGUAGE DETECTION SUMMARY",
        "─" * 80,
    ])
    ls = report.language_sources
    lines.append(f"  Languages detected from slide (parentheses): {ls.get('slide_parens', 0)}")
    lines.append(f"  Languages detected from slide (text):        {ls.get('slide_text', 0)}")
    lines.append(f"  Languages detected from filename:            {ls.get('filename', 0)}")
    lines.append(f"  Languages not detected:                      {ls.get('not_found', 0)}")
    lines.append("")
    
    # Issue summary
    lines.extend([
        "─" * 80,
        "ISSUE SUMMARY",
        "─" * 80,
        f"  Songs with year not on slide:    {len(report.songs_missing_year)}",
        f"  Songs missing language:          {len(report.songs_missing_language)}",
        f"  Songs with missing meanings:     {len(report.songs_with_missing_meanings)}",
        f"  Total stanzas missing meaning:   {sum(len(v) for v in report.songs_with_missing_meanings.values())}",
        f"  Files with errors:               {len(report.files_with_errors)}",
        "",
    ])
    
    # Detailed: Missing Year
    if report.songs_missing_year:
        lines.extend([
            "─" * 80,
            "SONGS WITH YEAR NOT FOUND ON SLIDE",
            "(year was inferred from folder name or default)",
            "─" * 80,
        ])
        for item in report.songs_missing_year:
            lines.append(f"  • {item}")
        lines.append("")
    
    # Detailed: Missing Language
    if report.songs_missing_language:
        lines.extend([
            "─" * 80,
            "SONGS MISSING LANGUAGE",
            "─" * 80,
        ])
        for item in report.songs_missing_language:
            lines.append(f"  • {item}")
        lines.append("")
    
    # Detailed: Missing Meanings
    if report.songs_with_missing_meanings:
        lines.extend([
            "─" * 80,
            "SONGS WITH MISSING ENGLISH MEANINGS",
            "(lists which slide numbers are missing translations)",
            "─" * 80,
        ])
        for song, slides in sorted(report.songs_with_missing_meanings.items()):
            lines.append(f"  • {song}")
            lines.append(f"      Missing on slides: {slides}")
        lines.append("")
    
    # Files with errors
    if report.files_with_errors:
        lines.extend([
            "─" * 80,
            "FILES WITH ERRORS",
            "─" * 80,
        ])
        for f in report.files_with_errors:
            lines.append(f"  • {f}")
        lines.append("")
    
    # Skipped files
    if report.skipped_files:
        lines.extend([
            "─" * 80,
            "SKIPPED FILES (no valid content extracted)",
            "─" * 80,
        ])
        for f in report.skipped_files:
            lines.append(f"  • {f}")
        lines.append("")
    
    # Full issue log
    lines.extend([
        "─" * 80,
        "DETAILED ISSUE LOG",
        "─" * 80,
        "",
    ])
    
    # Group by type
    by_type: dict[str, list[Issue]] = {}
    for issue in report.issues:
        t = issue.issue_type.name
        if t not in by_type:
            by_type[t] = []
        by_type[t].append(issue)
    
    for issue_type in sorted(by_type.keys()):
        issues = by_type[issue_type]
        lines.append(f"### {issue_type.replace('_', ' ')} ({len(issues)} occurrences)")
        lines.append("")
        
        for issue in issues:
            lines.append(f"  File: {issue.source_file}")
            if issue.full_path:
                lines.append(f"  Path: {issue.full_path}")
            if issue.song_title:
                lines.append(f"  Song: {issue.song_title}")
            if issue.slide_number:
                lines.append(f"  Slide: {issue.slide_number}")
            if issue.details:
                lines.append(f"  Details: {issue.details}")
            if issue.raw_text:
                lines.append(f"  Raw text sample:")
                for raw_line in issue.raw_text.split('\n')[:5]:
                    lines.append(f"    | {raw_line}")
            lines.append("")
    
    lines.extend([
        "=" * 80,
        "END OF REPORT",
        "=" * 80,
    ])
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    print(f"Report saved: {output_path}")
    
    # JSON report
    json_data = {
        'timestamp': report.timestamp,
        'input_path': report.input_path,
        'output_path': report.output_path,
        'summary': {
            'total_files_found': report.total_files_found,
            'total_files_processed': report.total_files_processed,
            'total_songs_extracted': report.total_songs_extracted,
            'total_stanzas': report.total_stanzas,
        },
        'songs_by_year': report.songs_by_year,
        'songs_missing_year': report.songs_missing_year,
        'songs_missing_language': report.songs_missing_language,
        'songs_with_missing_meanings': report.songs_with_missing_meanings,
        'files_with_errors': report.files_with_errors,
        'skipped_files': report.skipped_files,
        'issues': [issue.to_dict() for issue in report.issues],
    }
    
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(json_data, f, indent=2, ensure_ascii=False)
    print(f"JSON report saved: {json_path}")


def print_summary():
    """Print summary to console."""
    report = get_report()
    
    print(f"\n{'─' * 50}")
    print("SUMMARY")
    print(f"{'─' * 50}")
    print(f"  Files processed:    {report.total_files_processed}/{report.total_files_found}")
    print(f"  Songs extracted:    {report.total_songs_extracted}")
    print(f"  Total stanzas:      {report.total_stanzas}")
    
    if report.songs_missing_year:
        print(f"  ⚠ Year not on slide: {len(report.songs_missing_year)} songs")
    if report.songs_missing_language:
        print(f"  ⚠ Missing language:   {len(report.songs_missing_language)} songs")
    if report.songs_with_missing_meanings:
        total = sum(len(v) for v in report.songs_with_missing_meanings.values())
        print(f"  ⚠ Missing meanings:   {total} stanzas in {len(report.songs_with_missing_meanings)} songs")
    if report.files_with_errors:
        print(f"  ✗ Files with errors:  {len(report.files_with_errors)}")
    
    print(f"{'─' * 50}")
    print("  See report file for full details")


# =============================================================================
# MAIN
# =============================================================================

def run_pipeline(input_path: str, output_dir: str, output_name: str = None,
                years: list[str] = None, xml_only: bool = False, from_xml: str = None):
    """Run pipeline."""
    reset_report()
    report = get_report()
    
    output_name = output_name or DEFAULT_OUTPUT_NAME
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    
    xml_file = out / f"{output_name}{XML_EXTENSION}"
    pdf_file = out / f"{output_name}{PDF_EXTENSION}"
    toc_file = out / f"{output_name}{TOC_EXTENSION}"
    report_file = out / f"{output_name}{REPORT_EXTENSION}"
    report_json = out / f"{output_name}{REPORT_JSON_EXTENSION}"
    
    report.input_path = os.path.abspath(input_path)
    report.output_path = str(out.absolute())
    
    print("=" * 60)
    print("SONGBOOK GENERATOR")
    print("=" * 60)
    
    if from_xml:
        print(f"\n[1] Loading from XML: {from_xml}")
        songs = xml_to_songs(from_xml)
        print(f"  Loaded {len(songs)} songs")
    else:
        print(f"\n[1] Extracting from PPTX files")
        print(f"  Input: {input_path}")
        
        base = Path(input_path)
        years = years or YEAR_FOLDERS
        
        has_year_folders = any((base / y).exists() for y in years)
        
        if has_year_folders:
            songs = process_all_years(input_path, years)
        else:
            print(f"\n  Single folder mode (assuming year: {CURRENT_YEAR})")
            songs = process_year_folder(input_path, CURRENT_YEAR)
        
        if not songs:
            print("\n✗ No songs extracted!")
            generate_report(str(report_file), str(report_json))
            sys.exit(1)
        
        print(f"\n[2] Saving XML")
        songs_to_xml(songs, str(xml_file))
    
    if xml_only:
        print_summary()
        generate_report(str(report_file), str(report_json))
        print(f"\n{'=' * 60}")
        print("Done (XML only)")
        return
    
    step = 3 if not from_xml else 2
    print(f"\n[{step}] Generating PDF")
    songs = generate_pdf(songs, str(pdf_file))
    
    print(f"\n[{step+1}] Generating TOC")
    generate_toc(songs, str(toc_file))
    
    print(f"\n[{step+2}] Generating Report")
    generate_report(str(report_file), str(report_json))
    
    print_summary()
    
    print(f"\n{'=' * 60}")
    print("COMPLETE")
    print(f"{'=' * 60}")
    print(f"  XML:    {xml_file}")
    print(f"  PDF:    {pdf_file}")
    print(f"  TOC:    {toc_file}")
    print(f"  Report: {report_file}")


# =============================================================================
# TESTS & DEBUGGING
# =============================================================================

def test_language_detection():
    """Test suite for language detection logic."""
    print("\n" + "="*60)
    print("TESTING LANGUAGE DETECTION")
    print("="*60)
    
    test_cases = [
        # (slide_text, filename, expected_lang, expected_src)
        ("2025 (Malayalam)", "Song.pptx", "Malayalam", "slide_parens"),
        ("Malayalam\n2025", "Song.pptx", "Malayalam", "slide_text"),
        ("2025\n(mal)", "Song.pptx", "Malayalam", "slide_parens"),
        ("", "Prema_Sagarame_Mal_2025.pptx", "Malayalam", "filename"),
        ("", "Tamil_Song.pptx", "Tamil", "filename"),
        ("2025 (Tamil)", "Malayalam_Song.pptx", "Tamil", "slide_parens"), # Priority check
        ("Bhajan in Telugu", "Song.pptx", "Telugu", "slide_text"),
        ("New Song 2024 Kannada", "Song.pptx", "Kannada", "slide_text"),
        ("(Bengali) Title", "Song.pptx", "Bengali", "slide_parens"),
        ("", "Song (Mar).pptx", "Marathi", "filename"),
        ("(braj) Song", "Song.pptx", "Braj Bhasha", "slide_parens"),
        ("Song in Braja", "Song.pptx", "Braj Bhasha", "slide_text"),
        ("", "Song_mwr_2025.pptx", "Marwari", "filename"),
        ("Song (Marwadi)", "Song.pptx", "Marwari", "slide_parens"),
        ("Song (Odiya)", "Song.pptx", "Odia", "slide_parens"),
    ]
    
    passed = 0
    for slide, file, exp_lang, exp_src in test_cases:
        lang, src = detect_language(slide_text=slide, filename=file)
        
        status = "✓" if lang == exp_lang and src == exp_src else "✗"
        slide_display = slide.replace('\n', ' ')
        print(f"  {status} Input: slide='{slide_display}', file='{file}'")
        print(f"    Expected: {exp_lang} ({exp_src})")
        print(f"    Got:      {lang} ({src})")
        
        if status == "✓":
            passed += 1
    
    print(f"\nResult: {passed}/{len(test_cases)} passed")
    return passed == len(test_cases)


def test_line_break_preservation(sample_pptx_path: str):
    """Debug line break preservation for a specific file."""
    if not os.path.exists(sample_pptx_path):
        print(f"File not found: {sample_pptx_path}")
        return
    
    print("\n" + "="*60)
    print(f"DEBUGGING TEXT EXTRACTION: {os.path.basename(sample_pptx_path)}")
    print("="*60)
    
    try:
        prs = Presentation(sample_pptx_path)
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i+1} ---")
            shapes = get_shapes_by_position(slide)
            for s_idx, shape in enumerate(shapes):
                text = shape['text']
                # Visualize line breaks
                visible_breaks = text.replace('\n', '[\\n]\n')
                print(f"  Shape {s_idx+1}:")
                print(f"  {visible_breaks}")
    except Exception as e:
        print(f"Error: {e}")


def main():
    parser = argparse.ArgumentParser(description='Convert PPTX songs to PDF songbook')
    parser.add_argument('input_path', nargs='?', help='Directory with PPTX files or year folders')
    parser.add_argument('-o', '--output-dir', default='output')
    parser.add_argument('-n', '--name', default=DEFAULT_OUTPUT_NAME)
    parser.add_argument('--years', nargs='+')
    parser.add_argument('--xml-only', action='store_true')
    parser.add_argument('--from-xml', metavar='FILE')
    parser.add_argument('--test', action='store_true', help='Run language detection tests')
    parser.add_argument('--debug-file', metavar='FILE', help='Debug text extraction for a file')
    
    args = parser.parse_args()
    
    if args.test:
        test_language_detection()
        if not args.input_path:
            return

    if args.debug_file:
        test_line_break_preservation(args.debug_file)
        if not args.input_path:
            return
            
    if not args.input_path:
        parser.print_help()
        return

    run_pipeline(args.input_path, args.output_dir, args.name, args.years, args.xml_only, args.from_xml)


if __name__ == "__main__":
    main()