#!/usr/bin/env python3
"""
pptx_parser.py - PPTX parsing logic for songbook generator.
"""

import os
import traceback
from pathlib import Path
from typing import Optional

from pptx import Presentation

from config import CURRENT_YEAR, YEAR_FOLDERS
from models import Issue, IssueType, Song, Stanza, get_report
from language import detect_language, detect_year


def extract_text_from_shape(shape) -> str:
    """
    Extract text from PowerPoint shape, preserving paragraph breaks only.
    
    PowerPoint stores text in paragraphs. Each paragraph is a separate line.
    Within a paragraph, runs are joined together (they're just formatting spans).
    """
    if not shape.has_text_frame:
        return ""
    
    paragraphs = []
    
    for paragraph in shape.text_frame.paragraphs:
        # Collect runs and join them - runs are just formatting spans within a paragraph
        runs_texts = []
        
        for run in paragraph.runs:
            run_text = run.text
            # Handle various line break characters that might be in PPTX:
            run_text = run_text.replace('\r\n', '\n')  # CRLF first
            run_text = run_text.replace('\r', '\n')    # Then CR
            run_text = run_text.replace('\v', '\n')    # Vertical tab
            run_text = run_text.replace('\x0b', '\n')  # Same as \v
            run_text = run_text.replace('\u000B', '\n')  # Unicode vertical tab
            run_text = run_text.replace('\u2028', '\n')  # Line separator
            run_text = run_text.replace('\u2029', '\n')  # Paragraph separator
            if run_text:  # Only add non-empty runs
                runs_texts.append(run_text)
        
        # Join runs with empty string - they're just formatting spans, not separate lines
        para_text = ''.join(runs_texts).strip()
        paragraphs.append(para_text)
    
    # Join paragraphs with newlines
    result = '\n'.join(paragraphs)
    
    # Replace non-breaking spaces with regular spaces
    result = result.replace('\xa0', ' ').replace('\u00A0', ' ')
    
    # Clean up: remove leading/trailing blank lines but preserve internal structure
    lines = result.split('\n')
    
    # Strip leading empty lines
    while lines and not lines[0].strip():
        lines.pop(0)
    
    # Strip trailing empty lines  
    while lines and not lines[-1].strip():
        lines.pop()
    
    return '\n'.join(lines)




def get_shapes_by_position(slide) -> list[dict]:
    """Get text shapes sorted by vertical then horizontal position."""
    shapes = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text.strip():
                shapes.append({
                    'top': shape.top or 0,
                    'left': shape.left or 0,
                    'text': text,
                    'raw': text,
                })
    
    shapes.sort(key=lambda x: (x['top'], x['left']))
    return shapes


def parse_title_slide(
    slide, 
    source_file: str, 
    full_path: str,
    folder_year: Optional[str] = None
) -> tuple[Optional[str], Optional[str], Optional[str], str, str]:
    """
    Parse title slide for title, year, and language.
    
    Returns:
        (title, year, language, year_source, language_source)
    """
    report = get_report()
    shapes = get_shapes_by_position(slide)
    
    if not shapes:
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_TITLE,
            source_file=source_file,
            full_path=full_path,
            slide_number=1,
            details="No text found on title slide (slide 1)",
        ))
        return None, None, None, "", ""
    
    # Collect all text for analysis
    all_texts = [s['text'] for s in shapes]
    all_text_combined = "\n---\n".join(all_texts)
    
    # First shape is typically the title
    title_text = shapes[0]['text'].strip()
    
    # Handle multi-line titles
    if '\n' in title_text:
        lines = [l.strip() for l in title_text.split('\n') if l.strip()]
        title = lines[0] if lines else ""
    else:
        title = title_text
    
    if not title:
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_TITLE,
            source_file=source_file,
            full_path=full_path,
            slide_number=1,
            details="Title text box is empty",
            raw_text=all_text_combined[:500],
        ))
        return None, None, None, "", ""
    
    # --- YEAR DETECTION ---
    year = None
    year_source = ""
    
    for text in all_texts:
        found_year = detect_year(text)
        if found_year:
            year = found_year
            year_source = "slide"
            break
    
    if not year:
        found_year = detect_year(source_file)
        if found_year:
            year = found_year
            year_source = "filename"
    
    if not year and folder_year:
        year = folder_year
        year_source = "folder"
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_YEAR,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            slide_number=1,
            details=f"Year not found on slide or filename, using folder: {folder_year}",
            raw_text=all_text_combined[:300],
        ))
    
    if not year:
        year = CURRENT_YEAR
        year_source = "default"
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_YEAR,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            slide_number=1,
            details=f"Year not found anywhere, using default: {CURRENT_YEAR}",
            raw_text=all_text_combined[:300],
        ))
    
    # --- LANGUAGE DETECTION ---
    language, language_source = detect_language(
        slide_text=all_text_combined,
        filename=source_file,
        all_slide_texts=all_texts
    )
    
    if not language:
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_LANGUAGE,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            slide_number=1,
            details=(
                "Language not found. Checked:\n"
                f"  - Title slide text\n"
                f"  - Filename: {source_file}\n"
                "Expected formats: '(Malayalam)', '2025 Tamil', 'Song_Mal_2025.pptx'"
            ),
            raw_text=all_text_combined[:400],
        ))
    
    return title, year, language, year_source, language_source


def parse_lyric_slide(
    slide,
    slide_number: int,
    source_file: str,
    full_path: str,
    song_title: str,
    language: Optional[str] = None
) -> tuple[Optional[str], Optional[str], bool]:
    """
    Parse lyric slide.
    
    Returns:
        (lyrics, meaning, had_error)
    """
    report = get_report()
    
    try:
        shapes = get_shapes_by_position(slide)
    except Exception as e:
        report.add_issue(Issue(
            issue_type=IssueType.SLIDE_ERROR,
            source_file=source_file,
            full_path=full_path,
            song_title=song_title,
            slide_number=slide_number,
            details=f"Error reading slide: {str(e)}",
        ))
        return None, None, True
    
    if not shapes:
        report.add_issue(Issue(
            issue_type=IssueType.EMPTY_STANZA,
            source_file=source_file,
            full_path=full_path,
            song_title=song_title,
            slide_number=slide_number,
            details=f"Slide {slide_number} has no text content",
        ))
        return None, None, False
    
    if len(shapes) == 1:
        lyrics = shapes[0]['text']
        
        # English songs don't need translation/meaning
        if language != "English":
            report.add_issue(Issue(
                issue_type=IssueType.MISSING_MEANING,
                source_file=source_file,
                full_path=full_path,
                song_title=song_title,
                slide_number=slide_number,
                details=f"Slide {slide_number} has lyrics but no English meaning",
                raw_text=lyrics[:200] if lyrics else None,
            ))
        return lyrics, None, False
    
    # First shape = lyrics, last shape = meaning
    lyrics = shapes[0]['text']
    meaning = shapes[-1]['text']
    
    # Sanity check: if meaning looks like more lyrics (no English), flag it
    if language != "English" and meaning and not any(c.isascii() and c.isalpha() for c in meaning[:50]):
        report.add_issue(Issue(
            issue_type=IssueType.MISSING_MEANING,
            source_file=source_file,
            full_path=full_path,
            song_title=song_title,
            slide_number=slide_number,
            details=f"Slide {slide_number}: Bottom text doesn't appear to be English translation",
            raw_text=f"Top: {lyrics[:100]}...\nBottom: {meaning[:100]}...",
        ))
    
    return lyrics, meaning, False


def process_pptx_file(
    filepath: str, 
    folder_year: Optional[str] = None
) -> Optional[Song]:
    """Process a single PPTX file."""
    report = get_report()
    source_file = os.path.basename(filepath)
    full_path = os.path.abspath(filepath)
    
    try:
        prs = Presentation(filepath)
    except Exception as e:
        report.add_issue(Issue(
            issue_type=IssueType.FILE_ERROR,
            source_file=source_file,
            full_path=full_path,
            details=f"Cannot open file: {str(e)}\n{traceback.format_exc()}",
        ))
        return None
    
    slides = list(prs.slides)
    if not slides:
        report.add_issue(Issue(
            issue_type=IssueType.PARSE_ERROR,
            source_file=source_file,
            full_path=full_path,
            details="File contains no slides",
        ))
        return None
    
    report.total_slides_processed += len(slides)
    
    # Parse title slide
    title, year, language, year_source, language_source = parse_title_slide(
        slides[0], source_file, full_path, folder_year
    )
    
    if not title:
        return None
    
    # Parse lyric slides
    stanzas = []
    missing_meaning_slides = []
    
    for idx, slide in enumerate(slides[1:], start=2):
        lyrics, meaning, had_error = parse_lyric_slide(
            slide, idx, source_file, full_path, title, language
        )
        
        if had_error:
            continue
        
        if not lyrics:
            continue
        
        has_meaning = meaning is not None and meaning.strip() != ""
        
        if not has_meaning and language != "English":
            missing_meaning_slides.append(idx)
        
        stanzas.append(Stanza(
            lyrics=lyrics,
            meaning=meaning if has_meaning else None,
            slide_number=idx,
            has_meaning=has_meaning,
        ))
    
    if not stanzas:
        report.add_issue(Issue(
            issue_type=IssueType.NO_STANZAS,
            source_file=source_file,
            full_path=full_path,
            song_title=title,
            details=f"No valid lyric slides found. File has {len(slides)} slides total.",
        ))
    
    return Song(
        title=title,
        year=year,
        language=language,
        source_file=source_file,
        full_path=full_path,
        stanzas=stanzas,
        year_source=year_source,
        language_source=language_source,
        has_year=(year_source == "slide"),
        has_language=(language is not None),
        missing_meaning_slides=missing_meaning_slides,
    )


def process_year_folder(folder_path: str, year: str) -> list[Song]:
    """Process all PPTX in a year folder."""
    report = get_report()
    songs = []
    folder = Path(folder_path)
    
    if not folder.exists():
        print(f"  ⚠ Folder not found: {folder_path}")
        return songs
    
    pptx_files = [f for f in sorted(folder.glob("*.pptx")) if not f.name.startswith("~$")]
    
    if not pptx_files:
        print(f"  ⚠ No .pptx files in {folder_path}")
        return songs
    
    report.total_files_found += len(pptx_files)
    report.files_by_year[year] = len(pptx_files)
    
    print(f"  Found {len(pptx_files)} files")
    
    for filepath in pptx_files:
        song = process_pptx_file(str(filepath), folder_year=year)
        
        if song:
            songs.append(song)
            report.total_files_processed += 1
            report.total_songs_extracted += 1
            report.total_stanzas += len(song.stanzas)
            
            y = song.year or "Unknown"
            report.songs_by_year[y] = report.songs_by_year.get(y, 0) + 1
            
            # Update language source stats
            if not song.has_language:
                report.language_sources["not_found"] += 1
            else:
                src = song.language_source or "slide_text"
                report.language_sources[src] = report.language_sources.get(src, 0) + 1
            
            # Detailed status
            issues = []
            if not song.has_language:
                issues.append("no language detected")
            elif song.language_source == "filename":
                issues.append(f"language from filename: {song.language}")
            
            if song.year_source == "folder":
                issues.append(f"year from folder: {song.year}")
            elif song.year_source == "filename":
                issues.append(f"year from filename: {song.year}")
            elif song.year_source == "default":
                issues.append(f"year defaulted: {song.year}")
            
            if song.missing_meaning_slides:
                issues.append(f"missing meaning: slides {song.missing_meaning_slides}")
            
            lang_display = f" ({song.language})" if song.language else ""
            
            if issues:
                print(f"    ⚠ {song.title}{lang_display}")
                for issue in issues:
                    print(f"        └─ {issue}")
            else:
                print(f"    ✓ {song.title}{lang_display}")
        else:
            report.skipped_files.append(str(filepath))
            print(f"    ✗ {filepath.name}")
            print(f"        └─ skipped (see report for details)")
    
    return songs


def process_all_years(base_path: str, years: list[str] = None) -> list[Song]:
    """Process all year folders."""
    if years is None:
        years = YEAR_FOLDERS
    
    all_songs = []
    base = Path(base_path)
    
    for year in years:
        year_folder = base / year
        print(f"\n{'─' * 50}")
        print(f"Processing {year}/")
        print(f"{'─' * 50}")
        
        songs = process_year_folder(str(year_folder), year)
        all_songs.extend(songs)
        
        print(f"  → {len(songs)} songs from {year}")
    
    return all_songs
